"""Build the BookNest final presentation as an editable PowerPoint (.pptx),
following Final_Presentation_Outline.md and embedding the UI mockups."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
MOCK = os.path.join(HERE, "mockups")

GREEN = RGBColor(0x1B, 0x5E, 0x20)
DARK = RGBColor(0x20, 0x20, 0x20)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    return prs.slides.add_slide(BLANK)


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def set_run(p, text, size, color=DARK, bold=False, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return r


def title_bar(slide, text, subtitle=None):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]
    set_run(p, text, 30, WHITE, bold=True)
    if subtitle:
        p2 = tf.add_paragraph()
        set_run(p2, subtitle, 14, RGBColor(0xCF, 0xE0, 0xCF))
    return bar


def bullets(slide, items, left=Inches(0.6), top=Inches(1.5),
            width=Inches(7.6), height=Inches(5.4), size=18):
    tf = textbox(slide, left, top, width, height)
    first = True
    for item, lvl in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        prefix = "•  " if lvl == 0 else "–  "
        set_run(p, prefix + item, size - lvl * 2,
                DARK if lvl == 0 else GREY, bold=(lvl == 0 and item.endswith(":")))
        p.space_after = Pt(6)
    return tf


def add_mock(slide, name, left, top, height=Inches(5.6)):
    path = os.path.join(MOCK, name)
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, height=height)


# ---- Slide 1: Title ----
s = add_slide()
band = s.shapes.add_shape(1, 0, Inches(2.2), SW, Inches(2.9))
band.fill.solid(); band.fill.fore_color.rgb = GREEN; band.line.fill.background()
tf = band.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p, "BookNest", 54, WHITE, bold=True)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
set_run(p2, "Discover books recommended by students at your university", 20,
        RGBColor(0xE8, 0xF2, 0xE8), italic=True)
tf2 = textbox(s, Inches(0), Inches(5.4), SW, Inches(1.5))
for txt, sz, bold in [("Ozodbek Tursunaliyev  ·  Student ID 57264", 18, True),
                      ("Akademia WSB — Android + Firebase Team Project (individual)", 14, False)]:
    p = tf2.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    set_run(p, txt, sz, DARK, bold=bold)

# ---- Slide 2: Problem & users ----
s = add_slide(); title_bar(s, "Problem & Target Users")
bullets(s, [
    ("The problem:", 0),
    ("Students struggle to find relevant book recommendations from peers in similar subjects", 1),
    ("They waste time on generic lists that don't match academic interests", 1),
    ("Target users:", 0),
    ("University students (18–30), undergraduate and graduate", 1),
    ("Read for both academic enrichment and leisure", 1),
    ("Our solution:", 0),
    ("A peer-driven feed of book recommendations + a personal reading list", 1),
], width=Inches(12))

# ---- Slide 3: Solution overview (mockup) ----
s = add_slide(); title_bar(s, "Solution Overview")
bullets(s, [
    ("A social book-recommendation app for one campus", 0),
    ("Sign in, browse peer recommendations in real time", 0),
    ("Recommend books with a cover, genre and description", 0),
    ("Save books to a personal Want-to-Read list", 0),
    ("Built on Firebase (Auth, Firestore, Storage)", 0),
], width=Inches(7.4), size=19)
add_mock(s, "03_booklist.png", Inches(9.0), Inches(1.4))

# ---- Slide 4: Demo - accounts ----
s = add_slide(); title_bar(s, "Live Demo — Accounts",
                           "Register → land on the feed → log out → log back in (session persists)")
add_mock(s, "01_login.png", Inches(2.0), Inches(1.5), height=Inches(5.3))
add_mock(s, "02_register.png", Inches(7.4), Inches(1.5), height=Inches(5.3))

# ---- Slide 5: Demo - CRUD ----
s = add_slide(); title_bar(s, "Live Demo — Browse, Add, Edit, Delete",
                           "Real-time Firestore list; owner-only edit & delete with confirmation")
add_mock(s, "03_booklist.png", Inches(2.0), Inches(1.5), height=Inches(5.3))
add_mock(s, "04_detail.png", Inches(7.4), Inches(1.5), height=Inches(5.3))

# ---- Slide 6: Advanced feature ----
s = add_slide(); title_bar(s, "Advanced Feature — Firebase Storage",
                           "Pick a cover → upload with progress → stored URL shown via Coil")
bullets(s, [
    ("Why Storage:", 0),
    ("Cover images give instant visual recognition in the feed", 1),
    ("Flow:", 0),
    ("Gallery pick → preview → upload (progress bar) → URL saved to Firestore", 1),
    ("Cover deleted from Storage when the book is deleted", 1),
    ("Secured:", 0),
    ("Authenticated writes only, images < 5 MB (storage.rules)", 1),
], width=Inches(7.2), size=18)
add_mock(s, "05_addbook.png", Inches(9.0), Inches(1.4))

# ---- Slide 7: Architecture ----
s = add_slide(); title_bar(s, "Architecture")
bullets(s, [
    ("Android client — Kotlin + Jetpack Compose (Material 3), MVVM", 0),
    ("State via StateFlow; async Firebase calls via coroutines", 0),
    ("Firebase backend (serverless, Spark free tier):", 0),
    ("Authentication — email/password", 1),
    ("Cloud Firestore — users, books, savedBooks (real-time listeners)", 1),
    ("Firebase Storage — book_covers/", 1),
    ("Clean layering: UI → ViewModel → pure logic (util) → data models", 0),
    ("Pure logic is unit-tested on the JVM (no emulator needed)", 0),
], width=Inches(12), size=18)

# ---- Slide 8: Data model ----
s = add_slide(); title_bar(s, "Data Model — 3 Firestore Collections")
bullets(s, [
    ("users:", 0),
    ("userId, email, displayName, university, createdAt", 1),
    ("books:", 0),
    ("bookId, title, author, description, genre, coverImageUrl, recommendedBy, recommenderName, createdAt", 1),
    ("savedBooks:", 0),
    ("savedId, userId, bookId, title, author, status, savedAt", 1),
    ("Security rules:", 0),
    ("Authenticated access; only the recommender can edit/delete their books", 1),
], width=Inches(12), size=18)

# ---- Slide 9: Testing ----
s = add_slide(); title_bar(s, "Testing")
bullets(s, [
    ("17 JUnit unit tests, all passing (./gradlew test)", 0),
    ("ValidatorsTest — email, password, password-match, book-input rules", 1),
    ("BookFiltersTest — search by title/author/genre, sorting, combined", 1),
    ("Logic kept as pure functions (util/) → fast JVM tests, no emulator", 0),
    ("Manual end-to-end testing of login → CRUD → cover upload each lab", 0),
], width=Inches(12), size=19)

# ---- Slide 10: Usability + saved mockup ----
s = add_slide(); title_bar(s, "Usability Testing & Improvements")
bullets(s, [
    ("Structured think-aloud test (Lab 3 protocol)", 0),
    ("Fixed: upload progress bar (was silent)", 0),
    ("Fixed: add-book validation with inline errors", 0),
    ("Fixed: bookmark toggles filled/outline", 0),
    ("Fixed: friendly empty-state message", 0),
    ("Backlogged: genre dropdown", 0),
], width=Inches(7.4), size=19)
add_mock(s, "06_saved.png", Inches(9.0), Inches(1.4))

# ---- Slide 11: Challenges ----
s = add_slide(); title_bar(s, "Challenges & Solutions")
bullets(s, [
    ("Firestore 'permission denied' → wrote authenticated security rules", 0),
    ("Blank cover URL → await downloadUrl before writing the document", 0),
    ("Remote covers not loading → INTERNET permission + Coil", 0),
    ("Back button returned to login → popUpTo(0) clears the back stack", 0),
    ("Solo work, no peer review → unit tests + structured self-usability test", 0),
], width=Inches(12), size=19)

# ---- Slide 12: Results & reflection ----
s = add_slide(); title_bar(s, "Results & Reflection")
bullets(s, [
    ("All deliverables met across Labs 1–3, delivered solo", 0),
    ("Working auth, CRUD, 4 extensions, Storage feature, 17 tests", 0),
    ("What went well: clean layering made testing & extensions fast", 0),
    ("What was hard: Firebase rules and the first Compose setup", 0),
    ("Future: genre dropdown, anti-spam Cloud Function, profile pictures, UI tests", 0),
], width=Inches(12), size=19)

# ---- Slide 13: Thank you ----
s = add_slide()
band = s.shapes.add_shape(1, 0, Inches(2.6), SW, Inches(2.3))
band.fill.solid(); band.fill.fore_color.rgb = GREEN; band.line.fill.background()
tf = band.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p, "Thank you — Questions?", 40, WHITE, bold=True)
tf2 = textbox(s, 0, Inches(5.1), SW, Inches(1.6))
for txt in ["github.com/Ozodiy/booknest-app",
            "Trello board + full reports in /docs",
            "BookNest · Ozodbek Tursunaliyev · 57264"]:
    p = tf2.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    set_run(p, txt, 16, DARK)

out = os.path.join(HERE, "BookNest_Presentation.pptx")
prs.save(out)
print("wrote", out, "with", len(prs.slides._sldIdLst), "slides")
